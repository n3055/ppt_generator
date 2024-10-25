import google.generativeai as genai
import os
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
import streamlit as st
import json
from io import BytesIO
import requests
from PIL import Image
import random
load_dotenv()
pix_key = st.secrets["PIX_API_KEY"]
API_KEY = st.secrets["API_KEY"]
st.title("lame:blue[PPT]generator")
st.header(":red[M]:orange[V]:red[J] :red[Template]")
topic = st.text_input("Topic (Just write the topic name Ex:Badminton)")
slides = st.number_input("No. of Slides",value=0,step=1)
credits = st.text_area("Credits:","1MJ23CG0XX-Name1\n1MJ23CG0XX-Name2\n1MJ23CG0XX-Name3\n1MJ23CG0XX-Name4\n1MJ23CG0XX-Name5")
if st.button("GET PPT",type="primary"):
    with st.spinner("Generating a Lame PPT for you...."):
        genai.configure(api_key=os.environ["API_KEY"])
        model = genai.GenerativeModel("gemini-1.5-flash")
        response = model.generate_content("Create a ppt containing a total of "+str(slides+2)+" slides ,on topic "+str(topic)+"""
                                          in proper json format.
                                          **Instruction for slides format**
                                          First slide is always intro or title slide with heading and subheading.
                                          For first slide replace value of content with a subheading about the topic and replace heading with the topic provided.
                                          Both title and subtitle must be short as possible,
                                          also try to make heading of every slide short as well.
                                          Each slide should contain 4 or more points in the content
                                          Contents of each slide should be ****detailed******(i.e. should explain each point)
                                          in such a way that I don't have to explain the slides
                                          but just read the slides and people should follow.Don't exceed
                                          more than ***65 words per slide***.
                                          Also recommend relevant images to search on web for each slide.
                                          Make sure that the recommended image is unique for each slide.Don't add thank you or contact page or reference page.
                                          Also remember to generate **exact number** of slides as told.Remember Detailing is IMPORTANT!!!!
                                          **Instructions for output format**
                                          Print the json in one single line.
                                          Don't use code block.Key and value should be in double quotes.
                                          Don't miss commas and brackets delimeters and double quotes
                                          Use this JSON schema:
                                          slide = {'Heading': str, 'content': list[str],"image":str}
                                          Return: list[side]
                                          these error :
                                          ["JSONDecodeError: Expecting property name enclosed in double quotes: line 1 column 2032 (char 2031)",
                                          "JSONDecodeError: Extra data: line 3 column 1 (char 7188)"
                                          "Number of slides generated is less than the specified number",
                                          "Content of the slides are not detailed"]
                                          shall be avoided
                                          """)
        str_json = str(response.text)
        cont = json.loads(response.text)
##Get image urls from pixabay
        imgs_link = []
        for c in cont:
            inp = ''.join(['+' if char == ' ' else char for char in c["image"]])
            rq = requests.get("https://pixabay.com/api/?key="+pix_key+"&q="+inp+"&image_type=photo").json()
            imgs_link.append(rq["hits"][random.randint(0,10)]["webformatURL"])
##Creating our ppt
        ppt = Presentation()
        blank_slide_layout = ppt.slide_layouts[6]
        flag = True
        title_slide = ppt.slides.add_slide(blank_slide_layout)
        left = top = 0
        slide_width = ppt.slide_width
        slide_height = ppt.slide_height
        title_slide.shapes.add_picture("tbg.PNG", left, top, slide_width, slide_height)
        intro_left = Inches(1.354)
        intro_top = Inches(1.968)
        intro_width = Inches(8.716)
        intro_height = Inches(1.011)
        sub_left = Inches(1.27)
        sub_top = Inches(3.283)
        sub_width = Inches(5.83)
        sub_height = Inches(0.933)
        test_left = Inches(5.96)
        test_top = Inches(4.96)
        test_height = Inches(1.614)
        test_width = Inches(3.968)
        intro = title_slide.shapes.add_textbox(intro_left,intro_top,intro_width,intro_height)
        intro_tf = intro.text_frame 
        intro_tf.word_wrap = True
        intro_tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        intro_tf.shrink_text_on_overflow = True 
        intro_p = intro_tf.paragraphs[0]
        intro_p.text = cont[0]["Heading"]
        intro_p.font.size = Pt(38)
        intro_p.font.color.rgb = RGBColor(255, 255, 255)
        intro_p.font.name = "Bahnschrift SemiBold"
        sub = title_slide.shapes.add_textbox(sub_left,sub_top,sub_width,sub_height)
        sub_tf = sub.text_frame 
        sub_tf.word_wrap = True
        sub_tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        sub_tf.shrink_text_on_overflow = True 
        sub_p = sub_tf.paragraphs[0]
        for ct in cont[0]["content"]:
            sub_p.text = ct
        sub_p.font.size = Pt(25)
        sub_p.font.color.rgb = RGBColor(255, 255, 255)
        sub_p.font.name = "Times New Roman"
        test = title_slide.shapes.add_textbox(test_left,test_top,test_width,test_height)
        test_tf = test.text_frame 
        test_tf.word_wrap = True
        test_tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        test_tf.shrink_text_on_overflow = True 
        test_p = test_tf.paragraphs[0]
        test_p.text = credits
        test_p.font.size = Pt(18)
#test_p.font.color.rgb = RGBColor(255, 255, 255)
        test_p.font.name = "Times New Roman"
        ig = 1
        for c in cont:
            if flag:
                flag = False
                continue
            slide = ppt.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture("bgg.PNG", left, top, slide_width, slide_height)
            tt = Inches(0.149)
            tl = Inches(1.279)
            tw = Inches(8.157)
            th = Inches(0.7)
            txBox = slide.shapes.add_textbox(tl,tt,tw,th)
            tf = txBox.text_frame 
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            tf.shrink_text_on_overflow = True 
            p = tf.paragraphs[0]
            p.text = c["Heading"]
            p.font.size = Pt(36)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = "Bahnschrift SemiBold"
            cl = Inches(4.779)
            ct = Inches(1.988)
            cw = Inches(4.8976)
            ch = Inches(3.73622)
            t2 = slide.shapes.add_textbox(cl,ct,cw,ch)
            ctf = t2.text_frame
            ctf.word_wrap = True
            ctf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            ctf.shrink_text_on_overflow = True 
            flag1 = True
            p2 = ctf.paragraphs[0]
            p2.text = c["content"][0]
            p2.font.size = Pt(18)
            p2.font.name = "Times New Roman"
            for pts in c["content"]:
                if flag1:
                    flag1=False
                    continue
                p2 = ctf.add_paragraph()
                p2.text = pts
                p2.font.size = Pt(18)
                p2.font.name = "Times New Roman"
            il = Inches(0.248)
            it = Inches(1.787)
            iw = Inches(3.992)
            ih = Inches(4.185)
            response = requests.get(imgs_link[ig])
            img = Image.open(BytesIO(response.content))
            temp_image_path = "temp_image.jpg"
            img.save(temp_image_path)
            pic = slide.shapes.add_picture(temp_image_path,il,it,iw,ih)
            image_ratio = pic.width / pic.height
            area_ratio = iw / ih

            if image_ratio > area_ratio:
                pic.width = int(iw)
                pic.height = int(iw/image_ratio)
            else:
                pic.height = int(ih)
                pic.width = int(ih*image_ratio)

            pic.left = int(il+(iw-pic.width)/2)
            pic.top = int(it+(ih-pic.height)/2)
            ig+=1
        ppt_buffer = BytesIO()
        ppt.save(ppt_buffer)
        ppt_buffer.seek(0)
        st.download_button("Download ppt",data=ppt_buffer,file_name=str(topic)+".pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
footer = """<style>.footer {position: fixed;left: 0;bottom: 0;width: 100%;background-color: #000;color: #5f6061;text-align: center;}</style><div class='footer'><p>Made with ❤️ by RISHI</p></div>"""
st.markdown(footer, unsafe_allow_html=True)


